"""
Update V1_Presy.pptx to reflect what's been built since the original deck:
  - 7 pipeline stages (was 6)
  - 14 triggers (was 13) — added function_headroom from Partnering Matrix
  - EXL official product vocabulary (XTRAKTO.AI, EXELIA.AI, etc.) with outcome metrics
  - 5-area P&C value chain taxonomy
  - tam-leads Claude Skill
  - Updated rankings (Travelers 89.3, Allstate 64.5)

New slides added:
  - "What good looks like — the Travelers brief" (the killer artifact)
  - "The discovery — where the real TAM is hiding" (Partnering Matrix headroom)
  - "How EXL teams use this — no terminal, just ask" (Skill + demo transition)
  - "What's next" (roadmap)

Original preserved at V1_Presy_backup.pptx
"""
from pptx import Presentation
from pptx.util import Pt

from pathlib import Path
REPO = Path(__file__).resolve().parents[2]                      # repo root (code/scripts/<script>.py)
DECK = str(REPO / "demo" / "V1_Presy.pptx")

p = Presentation(DECK)


def set_title(slide, text):
    """Set the title placeholder text."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.is_placeholder \
                and shape.placeholder_format.idx == 0:
            shape.text_frame.text = text
            return


def set_content(slide, lines, font_size_pt=14):
    """Replace content placeholder with given lines as paragraphs.
    Each item in lines becomes its own paragraph."""
    target = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape.is_placeholder \
                and shape.placeholder_format.idx != 0:
            target = shape
            break
    if target is None:
        # fallback: any non-title text frame
        title_shape = slide.shapes.title if hasattr(slide.shapes, 'title') else None
        for shape in slide.shapes:
            if shape.has_text_frame and shape != title_shape:
                target = shape
                break
    if target is None:
        return

    tf = target.text_frame
    # First paragraph: replace text
    tf.text = lines[0] if lines else ""
    if tf.paragraphs[0].runs:
        for run in tf.paragraphs[0].runs:
            run.font.size = Pt(font_size_pt)
    # Subsequent paragraphs
    for line in lines[1:]:
        para = tf.add_paragraph()
        para.text = line
        for run in para.runs:
            run.font.size = Pt(font_size_pt)


# ============ UPDATE EXISTING SLIDES ============

# Slide 1: Title (just keep "TAM" — it works)
# No subtitle placeholder on the layout, so leave it as is.

# Slide 2: What this project is — light edits, sharpen
set_title(p.slides[1], "What this project is")
set_content(p.slides[1], [
    "EXL does back-office work for insurance companies — claims processing, underwriting paperwork, billing, customer service.",
    "Their competitors are big outsourcing firms: Cognizant, Genpact, WNS, TCS.",
    "",
    "We have 5 internal competitive-intelligence spreadsheets — 31 sheets, ~1,860 rows across 4 different file vintages — tracking which competitor is parked at each insurance customer, what their executives have publicly said, and where EXL has openings.",
    "",
    "Mental model: imagine building competitive-intel + lead-gen for AWS's enterprise sales team. \"Which AWS customers have an Azure contract expiring soon and have been making noise?\" — exact same shape, but for insurance services.",
], font_size_pt=16)

# Slide 3: One row example — keep as is, minor polish
set_title(p.slides[2], "What a single piece of intel looked like")
set_content(p.slides[2], [
    "One row from one of the 5 source spreadsheets (F5 / Competitor Analysis / Row 8):",
    "",
    "    Customer:          Travelers — top-10 US insurer, $27B in annual premiums",
    "    Competitor:        Cognizant — large outsourcing firm (~350K employees)",
    "    Contract status:   Up for renewal in < 6 months",
    "    Footprint:         450 people on claims processing for Travelers",
    "    Internal note:     \"Champion Challenger discussions in motion. Rohit has met Mojgan.\"",
    "",
    "Translation: a household-name insurance company is about to re-bid a 450-person contract. The incumbent is in trouble. EXL has a foot in the door at the customer's CIO.",
    "",
    "But that fact was buried alongside ~1,860 other rows across 31 sheets in 5 separate files. Nothing stitches them together.",
], font_size_pt=14)

# Slide 4: Stages — updated to 7 + new EXL vocabulary
set_title(p.slides[3], "Seven stages — plus EXL's vocabulary")
set_content(p.slides[3], [
    "1. Read every cell. Audit proves every row accounted for (1,862 source rows, 0 dropped silently).",
    "2. Clean duplicate names. 30+ name variants collapsed (\"CTS\" = \"Cognizant\" = \"Cognizant Technology Solutions\").",
    "3. Build one profile per customer. 236 ranked accounts.",
    "4. Detect buying signals. 14 deterministic rules. 60 firings across the corpus.",
    "5. Score 0–100. Account fit + signal strength + solution match + relationship warmth.",
    "6. Synthesize the brief. 20 polished one-pagers in pipeline/leads/.",
    "7. Strategic heatmap. White-space map + Partnering Matrix benchmarks.",
    "",
    "NEW since V1: we adopted EXL's official vocabulary from the Addressable Market Tracker deck.",
    "Lead briefs now name real EXL products (XTRAKTO.AI, EXELIA.AI, NerveHub, Paymentor,",
    "Subrosource, MedConnection, Digital Finance Suite) with their actual outcome metrics",
    "(\"25-30% cost-of-operations reduction\", etc.), and use the 5-area P&C value chain taxonomy.",
], font_size_pt=13)

# Slide 5 in source (will become slide 7 after reorder): Punchline — updated numbers
set_title(p.slides[4], "What came out — and three things we caught along the way")
set_content(p.slides[4], [
    "Travelers ranked #1 with score 89.3/100. Far ahead of #2 Allstate (64.5).",
    "",
    "9 different trigger types fired on Travelers simultaneously — all corroborating active",
    "shopping and incumbent friction. The pipeline-generated brief cites the same source rows",
    "as the lead we'd hand-built earlier. Independent validation that the system works.",
    "",
    "Three other things the pipeline caught:",
    "",
    "1) Three high-NWP Strategic accounts have ZERO competitive intel in the spreadsheets:",
    "   Allstate (rank #4), AXIS, Amer Family. Either greenfield or a data gap — either way,",
    "   that's a first action item.",
    "",
    "2) The same internal Capability Gap spreadsheet existed in 3 versions inside one file.",
    "   Pipeline flagged 2 as superseded, safe to delete. ~140 stale rows of free cleanup.",
    "",
    "3) A NaN bug almost generated 217 false-positive leads (cost_pressure trigger). The",
    "   deterministic + audited pipeline made it obvious — 4-line fix. A black-box LLM",
    "   pipeline would have shipped 217 wrong leads silently.",
], font_size_pt=13)

# Slide 6 in source (will become slide 9 after reorder): Files table — update entries
old_files_slide = p.slides[5]
table = None
for shape in old_files_slide.shapes:
    if shape.has_table:
        table = shape.table
        break

if table:
    # Replace each cell in existing rows
    new_files = [
        ("File", "What you'd see when you open it"),
        ("produced_data/pipeline/leads/01_travelers_group.md",
         "Travelers sales brief — 89.3/100, score breakdown, 17 triggers cited to source rows, EXL anchor products (XTRAKTO.AI + EXELIA.AI + NerveHub + Paymentor) with outcome metrics, CxO map, pitch hook, risks, next 3 steps."),
        ("produced_data/pipeline/WHITESPACE_MAP.xlsx",
         "7-sheet workbook — Master Heatmap with EXL's 5-area P&C taxonomy, Partnering Matrix benchmarks sorted by headroom, White Space Ranking, Capability Gaps."),
        ("produced_data/pipeline/LEADS_DIGEST.md",
         "Top 10 ranked accounts with click-through links to each polished brief."),
        ("pipeline/COVERAGE_AUDIT.md",
         "PASS verdict — every source row reconciled, 100% datetime coverage on all 1,076 output rows."),
    ]
    for r_idx in range(min(len(table.rows), len(new_files))):
        for c_idx in range(2):
            cell = table.cell(r_idx, c_idx)
            cell.text = new_files[r_idx][c_idx]
            # Try to keep small font
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(11)

# Update the title for the files slide (current slide index 5)
set_title(p.slides[5], "Files you can open right now")


# ============ NEW SLIDES (added at end, reordered below) ============

# NEW Slide A: The Travelers brief (the killer artifact)
new_a = p.slides.add_slide(p.slide_layouts[1])  # Title and Content
set_title(new_a, "What good looks like — the Travelers brief")
set_content(new_a, [
    "Score: 89.3/100   |   Confidence: Very High   |   17 triggers fired (9 unique types)",
    "",
    "Anchor products the lead recommends (with the outcome metrics from EXL's deck):",
    "    EXL XTRAKTO.AI       — Document AI                 — 25-30% cost-of-ops reduction",
    "    EXL EXELIA.AI        — Conversational AI            — 5-10 pt NPS lift",
    "    EXL NerveHub         — Workflow orchestration       — end-to-end visibility",
    "    EXL Paymentor        — Payment automation           — 100% leakage elimination",
    "",
    "Pitch hook the pipeline writes automatically:",
    "",
    "    \"Sequenced displacement vehicle for Cognizant's 1000+ FTE footprint as multiple",
    "    concurrent scopes approach renewal. Lead with EXL XTRAKTO.AI + EXL EXELIA.AI —",
    "    typical outcome: 25-30% reduction in cost of operations, 30% improvement in",
    "    cycle time. Anchor narrative on Travelers Group's stated operating-leverage",
    "    commitments, with EXL's solution mapping to their VOC already complete.\"",
    "",
    "Every claim cited: F5/Competitor Analysis/R8-R12, F1/Buyer Priorities/R7-R20, F1/Top Insurers and Brokers/R12.",
], font_size_pt=12)

# NEW Slide B: Partnering Matrix discovery
new_b = p.slides.add_slide(p.slide_layouts[1])
set_title(new_b, "The discovery — where the real TAM is hiding")
set_content(new_b, [
    "Reading the Addressable Market Tracker deck (slides 10-11) gave us EXL's official Partnering Matrix —",
    "the % of each function we typically retain when we're inside an account. Sorted by displacement headroom:",
    "",
    "    Claims FNOL:                  EXL retains 10-20%   →   85% headroom   (very high)",
    "    UW Risk Assessment:           EXL retains 10-20%   →   85% headroom   (very high)",
    "    Account Set Up & Clearance:   EXL retains 25-35%   →   70% headroom   (very high)",
    "    Policy Servicing:             EXL retains 30-40%   →   65% headroom   (high)",
    "    Subrogation:                  EXL retains 40-50%   →   55% headroom   (high)",
    "    UW Info Gathering:            EXL retains 85-90%   →   12% headroom   (EXL mature)",
    "",
    "What this means for Travelers:",
    "    Cognizant has 1,000 FTEs at Policy Servicing.",
    "    EXL typically retains 35% of that function.",
    "    → ~2,857 FTE TAM at this single function at Travelers (i.e. ~$5M+ of expansion beyond just displacing Cognizant).",
    "",
    "The pipeline now fires a function_headroom trigger automatically when a competitor is parked in a",
    "high-headroom function. Travelers +2 firings; Liberty Mutual, Swiss Re, AON, Nationwide all gain visibility.",
], font_size_pt=12)

# NEW Slide C: How EXL teams use it — Skill + demo transition
new_c = p.slides.add_slide(p.slide_layouts[1])
set_title(new_c, "How EXL teams use this — no terminal, just ask")
set_content(new_c, [
    "Built as a Claude Code Skill — tam-leads. Lives in .claude/skills/ of this repo.",
    "Any EXL teammate with a Claude subscription installs it in seconds (just clone the repo).",
    "",
    "Then they type questions in chat. Claude reads the skill, runs the right command,",
    "formats the output. No Python. No terminal. No training. Just chat.",
    "",
    "Examples (the four I'll demo live in 2 minutes):",
    "",
    "    →  \"List all the accounts we have\"",
    "    →  \"Give me a full brief on Travelers\"",
    "    →  \"Where's the biggest white space across the portfolio?\"",
    "    →  \"What about Allstate — what's going on there?\"",
    "",
    "Live demo →",
], font_size_pt=18)

# NEW Slide D: Roadmap
new_d = p.slides.add_slide(p.slide_layouts[1])
set_title(new_d, "What's next")
set_content(new_d, [
    "V2 — Publish as a Claude Code plugin.",
    "    Any EXL teammate installs in 10 seconds. No git clone, no setup.",
    "",
    "V3 — Pull in external data sources.",
    "    SEC EDGAR (free public filings)        — fills data gaps at Allstate, AXIS, Amer Family",
    "    Earnings call transcripts              — live CxO priorities replacing 2020 VOC data",
    "    NAIC quarterly filings (free)          — combined ratio + premium trends",
    "    AM Best ratings                        — financial strength signals",
    "",
    "V4 — AWS deployment, per the existing project plan (deck slide 5).",
    "    Phase 1 already specced: ~$23/month (SQS + Lambda + S3 + RDS + EventBridge).",
    "    Phase 2 adds Bedrock + Claude for signal extraction: $600–1,600/month.",
    "    Total Year 1: $300–$19,500 depending on phase coverage.",
    "",
    "V5 — Score-weight calibration. After first 10–20 leads convert to wins or losses,",
    "    tune the 4-component scoring formula (account fit / signal / solution / relationship).",
    "",
    "Asks for this room:",
    "  1. Sales/BD team to validate the top 10 leads look credible.",
    "  2. Approval for AWS Phase 1 (~$23/month) to start the deployment story.",
    "  3. Gaps in EXL product mapping — what did we miss in the vocabulary?",
], font_size_pt=13)


# ============ REORDER SLIDES ============
# Source order after inserts:
#   0: Title (kept)
#   1: What this project is (updated)
#   2: One row example (kept)
#   3: Seven stages (updated)
#   4: Punchline (updated; was 'old slide 5')
#   5: Files (updated; was 'old slide 6')
#   6: Travelers brief (NEW A)
#   7: Partnering Matrix (NEW B)
#   8: How EXL teams use it (NEW C)
#   9: Roadmap (NEW D)
#
# Desired order:
#   1. Title                                  -> idx 0
#   2. What this project is                   -> idx 1
#   3. One row example                        -> idx 2
#   4. Seven stages                           -> idx 3
#   5. Travelers brief                        -> idx 6  (NEW)
#   6. Partnering Matrix discovery            -> idx 7  (NEW)
#   7. Punchline                              -> idx 4
#   8. How EXL teams use it (live demo here)  -> idx 8  (NEW)
#   9. Roadmap                                -> idx 9  (NEW)
#  10. Files (reference)                      -> idx 5

new_order = [0, 1, 2, 3, 6, 7, 4, 8, 9, 5]

sldIdLst = p.slides._sldIdLst
slides_in_xml = list(sldIdLst)
for s in slides_in_xml:
    sldIdLst.remove(s)
for idx in new_order:
    sldIdLst.append(slides_in_xml[idx])


# ============ SAVE ============

p.save(DECK)
print("Saved updated V1_Presy.pptx")
print("Original backed up at V1_Presy_backup.pptx")
print()
print("Final slide order:")
final = Presentation(DECK)
for i, slide in enumerate(final.slides, 1):
    title = ""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            title = shape.text_frame.text.split('\n')[0]
            break
    print(f"  {i:>2}. {title[:80]}")

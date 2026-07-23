#!/usr/bin/env python3
"""
dump.py — universal document profiler (Claude's "eyes" during tam-ingest).

Deterministic, no LLM. Detects the file type and emits a JSON profile Claude reads
to AUTHOR a schema card.

Spreadsheets (.xlsx/.xlsm/.csv/.tsv) — one profile per sheet:
  - detected header row (not assumed to be row 1)
  - per column: inferred type, null %, distinct count, sample values, role hint
  - N sample data rows WITH their real source row numbers (for provenance)
  - anomaly flags: two-row header band, empty/scaffold sheet, near-duplicate sheet

Presentations (.pptx/.pptm) — one profile per DECK (slides are the grain):
  - per slide: title, text lines/bullets, any tables, speaker notes, image count,
    with the real 1-based SLIDE NUMBER (for provenance — a slide is retrievable)
  - per slide: date_candidates (dates found ON the slide) for per-datapoint temporality
  - deterministic entity_hints (frequency-ranked capitalized phrases + the slides
    they appear on) — hints Claude curates into the card's entities, never gospel

Images (.png/.jpg/.jpeg/.webp/.tif/.bmp/.gif) — one profile per IMAGE (the file is the unit):
  - dimensions/mode/format, OCR'd text (if a tesseract engine is available),
    metadata dates (EXIF / PNG text chunks), and date candidates from OCR/filename/metadata
  - Claude authors the card from this PLUS a direct visual read of the image

Every profile surfaces temporality (date_candidates / meta_dates) so each datapoint can be
dated and marked stated (literally in the source) vs inferred.

Usage:
  python3 code/tam/dump.py "<file>" [--sheet NAME] [--rows N] [--json OUT]
"""
import sys, os, json, re, argparse, hashlib
from pathlib import Path
from datetime import datetime, date

SAMPLE_ROWS = 6
HEADER_SCAN = 15          # how many top rows to consider as a possible header
IMAGE_EXT = (".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff", ".bmp", ".gif")


# ---------- temporality: deterministic date-candidate scanner (shared) ----------
# Finds date-like strings so Claude can assign a date to each datapoint and mark it
# stated (literally in the source) vs inferred. Domain-blind; Claude curates.
_MONTHS_RE = r"(?:jan(?:uary)?|feb(?:ruary)?|mar(?:ch)?|apr(?:il)?|may|jun(?:e)?|jul(?:y)?|" \
             r"aug(?:ust)?|sep(?:t)?(?:ember)?|oct(?:ober)?|nov(?:ember)?|dec(?:ember)?)"
_DATE_PATTERNS = [
    re.compile(r"\b\d{4}-\d{2}-\d{2}\b"),                                  # 2025-09-01
    re.compile(r"\b\d{1,2}[/-]\d{1,2}[/-]\d{2,4}\b"),                      # 9/1/25, 01-09-2025
    re.compile(rf"\b{_MONTHS_RE}\.?\s*'?\d{{2,4}}\b", re.I),               # Sep 2025, Sept'21
    re.compile(rf"\b\d{{1,2}}\s+{_MONTHS_RE}\.?\s*'?\d{{2,4}}\b", re.I),   # 3 Aug 2020
    re.compile(r"\b(?:FY|CY|Q[1-4])\s*['\-]?\s*(?:FY)?\s*\d{2,4}\b", re.I),# FY26, Q3 2025
    re.compile(r"\b\d{4}\s*Q[1-4]\b", re.I),                              # 2025 Q3
    re.compile(r"\b(?:19|20)\d{2}\b"),                                     # bare year 2020
]


def find_dates(*texts):
    """Return unique date-like strings found across the given texts, in first-seen order."""
    seen, out = set(), []
    for t in texts:
        if not t:
            continue
        for pat in _DATE_PATTERNS:
            for m in pat.finditer(str(t)):
                v = m.group(0).strip()
                key = v.lower()
                if key not in seen:
                    seen.add(key); out.append(v)
    return out


# ---------- type inference ----------
def _is_int(v):
    try:
        int(str(v).replace(",", "").strip()); return True
    except Exception:
        return False

def _is_float(v):
    try:
        float(str(v).replace(",", "").replace("$", "").replace("%", "").strip()); return True
    except Exception:
        return False

def infer_type(values, colname=""):
    vals = [v for v in values if v not in (None, "", " ")]
    if not vals:
        return "string"
    cn = colname.lower()
    if isinstance(vals[0], (datetime, date)) or any(isinstance(v, (datetime, date)) for v in vals):
        return "date"
    str_vals = [str(v).strip() for v in vals]
    if all(v.lower() in ("yes", "no", "y", "n", "true", "false") for v in str_vals):
        return "bool"
    if all("%" in v or (_is_float(v) and 0 <= (float(re.sub(r"[^0-9.\-]", "", v) or 0)) <= 1) for v in str_vals) and "%" in "".join(str_vals):
        return "percent"
    if all(_is_int(v) for v in str_vals):
        nums = [int(v.replace(",", "")) for v in str_vals]
        if all(1900 <= n <= 2100 for n in nums) and ("year" in cn or "yr" in cn):
            return "year"
        return "int"
    if all(_is_float(v) for v in str_vals):
        if any(s in cn for s in ("rev", "premium", "nwp", "gwp", "tcv", "budget", "$", "revenue", "spend")):
            return "currency_usd"
        return "float"
    if any(s in cn for s in ("company", "group", "insurer", "carrier", "client", "account", "broker")):
        return "name_company"
    if any(s in cn for s in ("ceo", "cfo", "coo", "cxo", "contact", "leader", "owner", "exec", "name")):
        return "name_person"
    if any(v.startswith(("http://", "https://", "www.", "file://")) for v in str_vals):
        return "url"
    distinct = len(set(str_vals))
    if distinct <= 15 and distinct < max(2, len(str_vals) * 0.6):
        return "enum"
    if max((len(v) for v in str_vals), default=0) > 60:
        return "freetext"
    return "string"


def role_hint(t, null_pct):
    if t in ("name_company",):
        return "entity_key"
    if t in ("int", "float", "currency_usd", "percent"):
        return "measure"
    if t in ("enum", "bool", "year", "date"):
        return "dimension"
    if t in ("freetext", "url", "name_person"):
        return "attribute"
    return "attribute"


# ---------- header detection ----------
def score_row(cells):
    non_empty = [c for c in cells if c not in (None, "", " ")]
    if not non_empty:
        return -1
    strings = [c for c in non_empty if isinstance(c, str) and not _is_float(c)]
    uniq = len(set(str(c).strip() for c in non_empty))
    return len(non_empty) * 1.0 + len(strings) * 0.5 + uniq * 0.5


def detect_header(grid):
    best_i, best_s = 0, -1
    for i in range(min(HEADER_SCAN, len(grid))):
        s = score_row(grid[i])
        # a header should be followed by a data row (not empty)
        nxt = score_row(grid[i + 1]) if i + 1 < len(grid) else -1
        if nxt <= 0:
            continue
        if s > best_s:
            best_s, best_i = s, i
    # two-row header band: next row also mostly string labels with high overlap of emptiness
    two_row = False
    if best_i + 1 < len(grid):
        nxt = grid[best_i + 1]
        nxt_strings = [c for c in nxt if isinstance(c, str) and not _is_float(c) and c not in (None, "", " ")]
        nxt_non_empty = [c for c in nxt if c not in (None, "", " ")]
        if nxt_non_empty and len(nxt_strings) >= 0.7 * len(nxt_non_empty) and len(nxt_non_empty) < 0.6 * len(grid[best_i]):
            two_row = True
    return best_i, two_row


# ---------- readers ----------
def read_xlsx(path):
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    sheets = {}
    for ws in wb.worksheets:
        grid = []
        for row in ws.iter_rows(values_only=True):
            grid.append(list(row))
        sheets[ws.title] = grid
    wb.close()
    return sheets

def read_csv(path):
    import csv
    with open(path, newline="", encoding="utf-8", errors="ignore") as f:
        grid = [row for row in csv.reader(f)]
    return {Path(path).stem: grid}


# ---------- presentation reader (shared with query.py for grounded retrieval) ----------
def read_pptx(path):
    """Extract one dict per slide from a .pptx/.pptm. Deterministic, order-preserving.
    Each slide dict: {slide (1-based), layout, title, bullets[], text, tables[[...]],
    notes, n_images, char_count}. Used both by profiling (below) and by query.py to
    pull a slide's real text at answer time, so retrieval stays grounded in the file."""
    from pptx import Presentation
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        PICTURE = MSO_SHAPE_TYPE.PICTURE
    except Exception:                       # keep working if the enum import shifts
        PICTURE = 13
    prs = Presentation(path)
    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        try:
            title = slide.shapes.title.text.strip() if slide.shapes.title else None
        except Exception:
            title = None
        bullets, tables, n_images = [], [], 0
        for shape in slide.shapes:
            try:
                if shape.shape_type == PICTURE:
                    n_images += 1
            except Exception:
                pass
            if getattr(shape, "has_table", False):
                rows = [[(c.text or "").strip() for c in r.cells] for r in shape.table.rows]
                tables.append(rows)
                for r in rows:                       # table text is slide text too
                    line = " | ".join(x for x in r if x)
                    if line:
                        bullets.append(line)
                continue
            if getattr(shape, "has_text_frame", False):
                for p in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in p.runs).strip() or (p.text or "").strip()
                    if line:
                        bullets.append(line)
        # de-dupe the title line if it also came through as the first bullet
        if title and bullets and bullets[0].strip() == title.strip():
            body = bullets[1:]
        else:
            body = bullets
        notes = ""
        try:
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            notes = ""
        text = "\n".join(body)               # body only; the title is carried separately
        try:
            layout = slide.slide_layout.name
        except Exception:
            layout = None
        slides.append({
            "slide": idx,
            "layout": layout,
            "title": title,
            "bullets": body,
            "text": text,
            "tables": tables,
            "notes": notes,
            "n_images": n_images,
            "char_count": len(text),
        })
    return slides


# ---------- profiling ----------
def profile_sheet(name, grid, max_rows=SAMPLE_ROWS):
    # trim fully-empty trailing rows
    while grid and all(c in (None, "", " ") for c in grid[-1]):
        grid.pop()
    non_empty_rows = sum(1 for r in grid if any(c not in (None, "", " ") for c in r))
    if non_empty_rows <= 1:
        return {"sheet": name, "role": "non_data", "warnings": ["empty or scaffold sheet"],
                "n_rows": non_empty_rows, "columns": []}

    header_i, two_row = detect_header(grid)
    headers = grid[header_i]
    # de-dup + fill blank headers
    cols, seen = [], {}
    for j, h in enumerate(headers):
        h = (str(h).strip() if h not in (None, "") else f"col{j}")
        if h in seen:
            seen[h] += 1; h = f"{h}.{seen[h]}"
        else:
            seen[h] = 0
        cols.append(h)

    data = grid[header_i + (2 if two_row else 1):]
    ncols = len(cols)
    col_profiles = []
    for j in range(ncols):
        column = [r[j] if j < len(r) else None for r in data]
        non_null = [v for v in column if v not in (None, "", " ")]
        t = infer_type(non_null, cols[j])
        distinct = len(set(str(v).strip() for v in non_null))
        null_pct = round(100 * (len(column) - len(non_null)) / max(1, len(column)))
        samples = []
        for v in non_null:
            sv = str(v).strip()
            if sv not in samples:
                samples.append(sv)
            if len(samples) >= 5:
                break
        col_profiles.append({
            "name": cols[j], "type": t, "role_hint": role_hint(t, null_pct),
            "null_pct": null_pct, "distinct": distinct, "samples": samples,
        })

    # sample rows with real (1-based) source row numbers
    sample_rows = []
    first_data_excel_row = header_i + (2 if two_row else 1) + 1  # 1-based
    for k, r in enumerate(data):
        if any(c not in (None, "", " ") for c in r):
            obj = {cols[j]: (str(r[j]).strip() if j < len(r) and r[j] not in (None, "") else None)
                   for j in range(ncols)}
            sample_rows.append({"_source_row": first_data_excel_row + k, **obj})
        if len(sample_rows) >= max_rows:
            break

    colset_hash = hashlib.md5("|".join(sorted(cols)).encode()).hexdigest()[:10]
    return {
        "sheet": name, "role": "data",
        "header_row": header_i + 1,            # 1-based, for the card
        "two_row_header": two_row,
        "n_data_rows": sum(1 for r in data if any(c not in (None, "", " ") for c in r)),
        "n_cols": ncols,
        "colset_hash": colset_hash,
        "columns": col_profiles,
        "sample_rows": sample_rows,
        "warnings": (["two-row header band"] if two_row else []) +
                    ([f"trailing space in sheet name '{name}'"] if name != name.strip() else []),
    }


# ---------- presentation entity hints (deterministic; Claude curates into entities[]) ----------
# Common non-entity words that survive capitalization (slide starts, section labels).
_ENTITY_STOP = {
    "the", "a", "an", "and", "or", "of", "for", "to", "in", "on", "with", "by", "at", "as",
    "our", "we", "us", "you", "your", "this", "that", "these", "those", "it", "is", "are",
    "agenda", "overview", "summary", "introduction", "conclusion", "appendix", "contents",
    "slide", "page", "title", "note", "notes", "figure", "table", "source", "confidential",
    "draft", "final", "yes", "no", "tbd", "na", "n/a", "q1", "q2", "q3", "q4",
}
_ENTITY_RE = re.compile(r"\b([A-Z][A-Za-z0-9&.\-]*(?:\s+[A-Z][A-Za-z0-9&.\-]*)*)\b")


def entity_hints(slides, top=40):
    """Frequency-ranked capitalized phrases and the slides they appear on. A HINT for
    Claude — deterministic, domain-blind, never authoritative. Claude decides which are
    real entities, normalizes them, and writes the 'why'."""
    hits = {}            # lower -> {"name": display, "count": n, "slides": set}
    for s in slides:
        # scan line-by-line so a phrase never spans a line/newline; title once
        lines = [s.get("title") or ""] + (s.get("bullets") or []) + [s.get("notes") or ""]
        for line in lines:
          for m in _ENTITY_RE.finditer(line):
            phrase = m.group(1).strip(" .-&")
            if not phrase:
                continue
            words = phrase.split()
            # drop single short/stopword tokens; keep multi-word or ALLCAPS acronyms
            low = phrase.lower()
            single = len(words) == 1
            if single and (low in _ENTITY_STOP or (len(phrase) < 3) or
                           (not phrase.isupper() and len(phrase) < 4)):
                continue
            if all(w.lower() in _ENTITY_STOP for w in words):
                continue
            rec = hits.setdefault(low, {"name": phrase, "count": 0, "slides": set()})
            rec["count"] += 1
            rec["slides"].add(s["slide"])
    ranked = sorted(hits.values(), key=lambda r: (-r["count"], -len(r["slides"]), r["name"].lower()))
    return [{"name": r["name"], "count": r["count"], "slides": sorted(r["slides"])}
            for r in ranked[:top]]


# ---------- image reader / profiler (OCR + metadata) ----------
_EXIF_DATE_TAGS = {306: "DateTime", 36867: "DateTimeOriginal", 36868: "DateTimeDigitized"}


def read_image(path):
    """Extract deterministic material from a raster image: dimensions, embedded
    metadata dates (EXIF / PNG text chunks), and OCR'd text if a tesseract engine is
    available. Metadata dates are STATED temporality; OCR'd dates are content dates that
    may be stated (a label on the image) — Claude decides. Vision-level understanding is
    left to Claude reading the image directly (see the tam-ingest image track)."""
    from PIL import Image
    info = {"file": path, "kind": "image", "ocr_available": False, "ocr_text": None,
            "meta_dates": [], "warnings": []}
    try:
        img = Image.open(path)
    except Exception as e:
        info["warnings"].append(f"could not open image: {e}")
        return info
    info["width"], info["height"] = img.size
    info["mode"], info["format"] = img.mode, img.format

    # EXIF dates (JPEG/TIFF) — stated capture time in the file's metadata
    try:
        exif = img.getexif()
        for tag, name in _EXIF_DATE_TAGS.items():
            v = exif.get(tag)
            if v:
                info["meta_dates"].append({"source": f"exif:{name}", "value": str(v).strip()})
    except Exception:
        pass
    # PNG text chunks (e.g. 'Creation Time', 'date') — stated in the file's metadata
    try:
        for k, v in (getattr(img, "text", {}) or {}).items():
            if any(d in k.lower() for d in ("date", "time", "created")):
                info["meta_dates"].append({"source": f"png:{k}", "value": str(v).strip()})
    except Exception:
        pass

    # OCR (deterministic best-effort; engine optional)
    try:
        import pytesseract
        info["ocr_text"] = (pytesseract.image_to_string(img) or "").strip() or None
        info["ocr_available"] = True
    except Exception as e:
        info["warnings"].append(f"OCR engine unavailable ({e.__class__.__name__}); "
                                "author the transcript/summary by reading the image directly")
    return info


def profile_image(path):
    """Image profile: metadata + OCR text + date candidates. Slides/rows have a natural
    provenance unit; an image is one unit (the file). Claude authors an image card from
    this PLUS a direct visual read of the image."""
    info = read_image(path)
    ocr = info.get("ocr_text") or ""
    meta_date_vals = [d["value"] for d in info.get("meta_dates", [])]
    info["content_date_candidates"] = find_dates(ocr)                 # dates OCR'd from the image
    info["filename_date_candidates"] = find_dates(Path(path).name)    # a date in the filename
    info["meta_date_candidates"] = find_dates(*meta_date_vals)         # dates from file metadata
    if not ocr and info.get("ocr_available"):
        info["warnings"].append("OCR produced no text (photo/diagram?); rely on a direct visual read")
    return info


def profile_pptx(path, max_rows=SAMPLE_ROWS):
    """Deck-level profile: the collective raw material Claude needs to author a deck card
    (collective summary + per-slide point + curated entities). Slides are the grain."""
    slides = read_pptx(path)
    total_chars = sum(s["char_count"] for s in slides)
    slim = []
    for s in slides:
        # temporality: dates literally on/in this slide (a stated date for its datapoints)
        dates = find_dates(s["title"], s["text"], s["notes"])
        slim.append({
            "slide": s["slide"],
            "layout": s["layout"],
            "title": s["title"],
            "bullets": s["bullets"],
            "n_tables": len(s["tables"]),
            "tables": s["tables"][:2],          # first couple, for shape; full text is in bullets
            "has_notes": bool(s["notes"]),
            "notes": s["notes"],
            "n_images": s["n_images"],
            "char_count": s["char_count"],
            "date_candidates": dates,           # dates found ON this slide (stated); [] -> infer/inherit
        })
    return {
        "file": path,
        "kind": "presentation",
        "n_slides": len(slides),
        "total_chars": total_chars,
        "empty_slides": [s["slide"] for s in slides if s["char_count"] == 0 and s["n_images"] == 0],
        "filename_date_candidates": find_dates(Path(path).name),   # a deck-level stated date
        "slides": slim,
        "entity_hints": entity_hints(slides),
        "warnings": (["no extractable text (image-only deck?)"] if total_chars == 0 else []),
    }


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("file")
    ap.add_argument("--sheet")
    ap.add_argument("--rows", type=int, default=SAMPLE_ROWS)
    ap.add_argument("--json")
    args = ap.parse_args()
    path = args.file
    ext = Path(path).suffix.lower()

    # ----- presentations: deck profile (slides are the grain, not sheets) -----
    if ext in (".pptx", ".pptm"):
        out = profile_pptx(path, max_rows=args.rows)
        text = json.dumps(out, indent=2, default=str)
        if args.json:
            Path(args.json).parent.mkdir(parents=True, exist_ok=True)
            Path(args.json).write_text(text)
            print(f"Wrote {args.json}  ({out['n_slides']} slides)")
        else:
            print(text)
        return

    # ----- images: OCR + metadata profile (the file is one unit) -----
    if ext in IMAGE_EXT:
        out = profile_image(path)
        text = json.dumps(out, indent=2, default=str)
        if args.json:
            Path(args.json).parent.mkdir(parents=True, exist_ok=True)
            Path(args.json).write_text(text)
            print(f"Wrote {args.json}  (image {out.get('width')}x{out.get('height')}, "
                  f"ocr={'yes' if out.get('ocr_text') else 'no'})")
        else:
            print(text)
        return

    if ext in (".xlsx", ".xlsm"):
        sheets = read_xlsx(path)
    elif ext in (".csv", ".tsv"):
        sheets = read_csv(path)
    else:
        print(f"Unsupported: {ext} (spreadsheets: xlsx/xlsm/csv/tsv; presentations: pptx/pptm; "
              f"images: {'/'.join(e[1:] for e in IMAGE_EXT)})", file=sys.stderr)
        sys.exit(2)

    profiles = []
    for name, grid in sheets.items():
        if args.sheet and name.strip() != args.sheet.strip():
            continue
        profiles.append(profile_sheet(name, grid, max_rows=args.rows))

    # near-duplicate detection across sheets (same column set)
    by_hash = {}
    for p in profiles:
        h = p.get("colset_hash")
        if h:
            by_hash.setdefault(h, []).append(p["sheet"])
    for p in profiles:
        dupes = [s for s in by_hash.get(p.get("colset_hash"), []) if s != p["sheet"]]
        if dupes:
            p.setdefault("warnings", []).append(f"near-duplicate column set with: {dupes}")

    out = {"file": path, "n_sheets": len(profiles), "sheets": profiles}
    text = json.dumps(out, indent=2, default=str)
    if args.json:
        Path(args.json).write_text(text)
        print(f"Wrote {args.json}  ({len(profiles)} sheets)")
    else:
        print(text)


if __name__ == "__main__":
    main()
